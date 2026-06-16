"""Shared PPTX layout helpers for the Jieceng presentation decks.

Clean, emoji-free, 16:9 native slides built with python-pptx. Chinese text is
rendered with a CJK sans (Microsoft JhengHei, substituted on non-Windows hosts);
Latin/numerals use Inter / Consolas.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU = Inches
W = Inches(13.33)
H = Inches(7.5)

SANS = "Inter"
EA   = "Microsoft JhengHei"   # CJK; PowerPoint substitutes a local sans elsewhere
MONO = "Consolas"

WARM   = "#fafaf9"
CHAR   = "#1a1a1a"
BODY   = "#44403c"
STONE4 = "#a8a29e"
LINE   = "#e7e5e4"
WHITE  = "#ffffff"


def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb(color)


def _apply_font(run, size, bold, color, font, italic):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', EA if tag == 'a:ea' else font)


def text(slide, s, x, y, w, h, size=16, bold=False, color=CHAR,
         align=PP_ALIGN.LEFT, font=SANS, italic=False, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(s.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        _apply_font(r, size, bold, color, font, italic)
    return box


def bullets(slide, items, x, y, w, h, size=13, color=BODY, gap=6,
            marker_color=None, spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        if marker_color:
            rm = p.add_run()
            rm.text = "•  "
            _apply_font(rm, size, False, marker_color, SANS, False)
        r = p.add_run()
        r.text = it
        _apply_font(r, size, False, color, SANS, False)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def header(slide, kicker, title, accent, title_size=30):
    """Standard content-slide header: kicker + short accent rule + title."""
    text(slide, kicker, Inches(0.9), Inches(0.58), Inches(8), Inches(0.3),
         size=10.5, color=STONE4, font=MONO)
    rect(slide, Inches(0.9), Inches(1.0), Inches(0.55), Pt(2.4), fill=accent)
    text(slide, title, Inches(0.9), Inches(1.12), Inches(11.5), Inches(0.7),
         size=title_size, bold=True, color=CHAR)


def lead(slide, s, y=Inches(2.0), size=15, w=Inches(11.0)):
    return text(slide, s, Inches(0.9), y, w, Inches(1.2),
                size=size, color="#292524", spacing=1.25)


def note(slide, s, y=Inches(6.7)):
    return text(slide, s, Inches(0.9), y, Inches(11.5), Inches(0.6),
                size=10.5, color=STONE4, spacing=1.2)


def card(slide, x, y, w, h, accent=None, head=None, head_color=None,
         body=None, idx=None, fill=WHITE, body_items=None):
    rect(slide, x, y, w, h, fill=fill, line=LINE, line_w=0.75)
    if accent:
        rect(slide, x, y, w, Pt(2.6), fill=accent)
    pad = Inches(0.22)
    cy = y + Inches(0.18)
    if idx:
        text(slide, idx, x + pad, cy, w - pad, Inches(0.25), size=8.5,
             color=STONE4, font=MONO)
        cy += Inches(0.3)
    if head:
        text(slide, head, x + pad, cy, w - pad, Inches(0.4), size=14,
             bold=True, color=head_color or accent or CHAR)
        cy += Inches(0.5)
    if body:
        text(slide, body, x + pad, cy, w - pad, h - (cy - y) - Inches(0.1),
             size=11, color=BODY, spacing=1.2)
    if body_items:
        bullets(slide, body_items, x + pad, cy, w - pad,
                h - (cy - y) - Inches(0.1), size=10.5, color=BODY,
                marker_color=accent, gap=3, spacing=1.12)


def card_row(slide, items, y, h, cols=None, gap=Inches(0.28),
             left=Inches(0.9), total=Inches(11.53)):
    n = len(items)
    cw = (total - gap * (n - 1)) / n
    for i, it in enumerate(items):
        x = left + (cw + gap) * i
        card(slide, x, y, cw, h, **it)


def stat(slide, x, y, w, num, cap, color, num_size=30):
    text(slide, num, x, y, w, Inches(0.7), size=num_size, bold=True,
         color=color, font=SANS)
    text(slide, cap, x, y + Inches(0.62), w, Inches(0.3), size=11,
         color=STONE4)


def divider(prs, bg, kicker, title, subtitle, accent_kicker, accent_sub,
            footer="傑丞建築機構　Jieceng Architecture"):
    sl = blank(prs)
    fill_bg(sl, bg)
    text(sl, kicker, Inches(0.95), Inches(2.3), Inches(10), Inches(0.35),
         size=11, color=accent_kicker, font=MONO)
    text(sl, title, Inches(0.95), Inches(2.85), Inches(11), Inches(1.4),
         size=40, bold=True, color=WHITE, spacing=1.1)
    rect(sl, Inches(0.95), Inches(4.25), Inches(0.6), Pt(2.4), fill=accent_sub)
    text(sl, subtitle, Inches(0.95), Inches(4.5), Inches(10.5), Inches(0.9),
         size=16, color=accent_sub, spacing=1.3)
    text(sl, footer, Inches(0.95), Inches(6.7), Inches(10), Inches(0.3),
         size=11, color="#78716c", font=MONO)
    return sl


def cover(prs, bg, kicker, title, accent_lg, subtitle,
          footer="傑丞建築機構　Jieceng Architecture"):
    sl = blank(prs)
    fill_bg(sl, bg)
    text(sl, kicker, Inches(0.95), Inches(1.7), Inches(11), Inches(0.35),
         size=11, color=accent_lg, font=MONO)
    text(sl, title, Inches(0.95), Inches(2.25), Inches(11.3), Inches(2.0),
         size=46, bold=True, color=WHITE, spacing=1.12)
    rect(sl, Inches(0.95), Inches(4.55), Inches(0.6), Pt(2.4), fill=accent_lg)
    text(sl, subtitle, Inches(0.95), Inches(4.85), Inches(10.5), Inches(1.0),
         size=15, color="#a8a29e", spacing=1.35)
    text(sl, footer, Inches(0.95), Inches(6.75), Inches(10), Inches(0.3),
         size=11.5, color="#78716c", font=MONO)
    return sl


def closing(prs, bg, kicker, title, accent_lg,
            footer="傑丞建築機構　Jieceng Architecture"):
    sl = blank(prs)
    fill_bg(sl, bg)
    text(sl, kicker, Inches(0.95), Inches(2.4), Inches(11), Inches(0.35),
         size=11, color=accent_lg, font=MONO)
    text(sl, title, Inches(0.95), Inches(2.95), Inches(11), Inches(1.6),
         size=38, bold=True, color=WHITE, spacing=1.18)
    rect(sl, Inches(0.95), Inches(4.7), Inches(0.6), Pt(2.4), fill=accent_lg)
    text(sl, footer, Inches(0.95), Inches(6.7), Inches(10), Inches(0.3),
         size=11.5, color="#78716c", font=MONO)
    return sl
