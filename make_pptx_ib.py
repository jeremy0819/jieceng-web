"""Generate 智慧建築 IB PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)

def txb(slide, text, left, top, width, height,
        size=18, bold=False, color='#1a1a1a', align=PP_ALIGN.LEFT,
        italic=False):
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return tf_box

def txb_lines(slide, lines, left, top, width, height,
              size=14, color='#1a1a1a', bullet=False):
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ('• ' if bullet else '') + line
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return tf_box

def rect(slide, left, top, width, height, hex_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hex_color)
    shape.line.fill.background()
    return shape

def accent_bar(slide, left, top, width, hex_color):
    rect(slide, left, top, width, Pt(3), hex_color)

# Colors
BLUE_DARK  = '#0f172a'
BLUE       = '#1e3a8a'
BLUE_L     = '#1d4ed8'
BLUE_LG    = '#60a5fa'
CHARCOAL   = '#1a1a1a'
WHITE      = '#ffffff'
GRAY       = '#6b7280'
BLUE_BG    = '#eff6ff'
BLUE_XPALE = '#dbeafe'

C_ICT  = '#1d4ed8'  # blue
C_SEC  = '#dc2626'  # red
C_HLTH = '#16a34a'  # green
C_NRG  = '#d97706'  # amber
C_INT  = '#7c3aed'  # violet
C_CVN  = '#0891b2'  # cyan

def build_ib(path):
    prs = new_prs()

    # ── 01 Cover ──────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, BLUE_DARK)
    txb(sl,'內政部建築研究所 / Architecture and Building Research Institute',
        Inches(1),Inches(.6),Inches(11),Inches(.4),size=13,color=BLUE_LG)
    txb(sl,'智慧建築標章',
        Inches(1),Inches(1.2),Inches(11),Inches(1.0),size=48,bold=True,color=WHITE)
    txb(sl,'2024 版解析',
        Inches(1),Inches(2.15),Inches(11),Inches(.9),size=40,bold=True,color=BLUE_LG)
    rect(sl,Inches(1),Inches(3.2),Inches(.7),Pt(4),BLUE_LG)
    txb(sl,'六大評估指標 × 四個認證等級 × 2024 最新修訂重點',
        Inches(1),Inches(3.4),Inches(11),Inches(.4),size=16,color='#94a3b8')
    txb(sl,'從 ICT 基礎到 AI 數位孿生，全面解析台灣智慧建築認證制度',
        Inches(1),Inches(3.82),Inches(11),Inches(.4),size=16,color='#94a3b8')
    txb(sl,'2026.05.26',Inches(1),Inches(6.8),Inches(4),Inches(.4),size=12,color='#334155')

    # ── 02 目錄 ────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    accent_bar(sl,Inches(1),Inches(.8),Inches(1.2),BLUE_L)
    txb(sl,'今日議程',Inches(1),Inches(1.0),Inches(11),Inches(.6),size=32,bold=True,color=CHARCOAL)
    txb(sl,'PART 01 — 制度概覽',Inches(1),Inches(1.85),Inches(5.5),Inches(.35),size=13,color=BLUE_L)
    txb_lines(sl,['制度背景與沿革','2024 版六大更新重點','六大評估指標總覽','各指標配分比重'],
        Inches(1),Inches(2.25),Inches(5.5),Inches(1.8),size=15,color=CHARCOAL,bullet=True)
    txb(sl,'PART 02 — 指標深度解析',Inches(7),Inches(1.85),Inches(5.5),Inches(.35),size=13,color=BLUE_L)
    txb_lines(sl,['① 資訊通信','② 安全防災','③ 健康舒適','④ 節能管理','⑤ 設備整合','⑥ 貼心便利'],
        Inches(7),Inches(2.25),Inches(5.5),Inches(2.4),size=15,color=CHARCOAL,bullet=True)
    txb(sl,'PART 03 — 認證與申請',Inches(1),Inches(4.8),Inches(5.5),Inches(.35),size=13,color=BLUE_L)
    txb_lines(sl,['四個認證等級','申請流程說明'],
        Inches(1),Inches(5.2),Inches(5.5),Inches(.9),size=15,color=CHARCOAL,bullet=True)

    # ── 03 制度背景 ────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 制度背景',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'從 2004 到 2024：二十年智慧進化',
        Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    years=[('2004','智慧建築標章創立\n六大指標確立'),
           ('2010','修訂版本\n強化設備整合評估'),
           ('2016','IoT 概念\n納入評分'),
           ('2020','5G 時代\n前瞻條款 AI 準備'),
           ('2024★','AI × 數位孿生\n資安強化本次重點')]
    for i,(yr,desc) in enumerate(years):
        x=Inches(1+i*2.4)
        c=BLUE_LG if '★' in yr else '#94a3b8'
        rect(sl,x,Inches(2.0),Inches(.6),Inches(.6),c)
        txb(sl,yr.replace('★',''),x,Inches(2.0),Inches(.6),Inches(.6),
            size=11,bold=True,color=BLUE_DARK if '★' not in yr else CHARCOAL,align=PP_ALIGN.CENTER)
        txb(sl,yr,x-Inches(.1),Inches(2.72),Inches(.9),Inches(.3),
            size=11,bold=('★' in yr),color=BLUE_LG if '★' in yr else CHARCOAL)
        txb(sl,desc,x-Inches(.15),Inches(3.05),Inches(.95),Inches(.9),size=10,color=GRAY)
    txb(sl,'制度定位',Inches(1),Inches(4.1),Inches(6),Inches(.35),size=14,bold=True,color=CHARCOAL)
    txb(sl,'智慧建築標章由內政部建築研究所主管，評估建築物將資訊通信、節能管理、安全防災等\n六大智慧化系統整合至空間設計的程度，目標打造更安全、舒適、節能的人性化智慧建築。',
        Inches(1),Inches(4.55),Inches(7.5),Inches(1.0),size=13,color=CHARCOAL)
    for i,(num,lbl,bg) in enumerate([('6','評估指標','#f9fafb'),('4','認證等級','#f9fafb'),
                                      ('5yr','標章有效期','#f9fafb'),('ABRI','建築研究所主管',BLUE_BG)]):
        bx=Inches(9.0+i%2*1.9); by=Inches(4.1+i//2*1.7)
        rect(sl,bx,by,Inches(1.7),Inches(1.5),bg)
        txb(sl,num,bx,by+Inches(.1),Inches(1.7),Inches(.7),size=24,bold=True,color=BLUE_L,align=PP_ALIGN.CENTER)
        txb(sl,lbl,bx,by+Inches(.85),Inches(1.7),Inches(.4),size=10,color=GRAY,align=PP_ALIGN.CENTER)

    # ── 04 2024 更新重點 ────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'2024 版 ／ 修訂重點',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'2024 版六大更新重點',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    updates=[
        ('① ICT 更新','5G / Wi-Fi 6E 納入評分\n5G 小基站室內部署及 Wi-Fi 6E 全棟無縫覆蓋，正式列入評分。'),
        ('② 資安強化','零信任資安架構\n資訊安全防護獨立為必要評估條款，要求 IDS/IPS、VPN、Zero Trust 架構。'),
        ('③ AI 智慧化','各指標 AI 加分項\n六大指標均新增 AI 智慧化加分子項，包含 AI 影像分析、AI 需量預測等。'),
        ('④ 數位孿生','BIM 數位孿生整合\n設備整合指標新增 BIM 數位孿生加分，建築 3D 模型與 BMS 即時連動。'),
        ('⑤ 碳管理','建築碳足跡智慧監控\n節能管理指標新增碳排放即時監控，配合 2050 淨零路徑讓建築能耗透明化。'),
        ('⑥ 數位服務','行動 App 一站式服務\n貼心便利指標提升行動裝置整合評分，涵蓋物業 App、IoT 智慧家電整合。'),
    ]
    for i,(badge,body) in enumerate(updates):
        col=i%3; row=i//3
        bx=Inches(1+col*4.1); by=Inches(1.85+row*2.55)
        rect(sl,bx,by,Inches(3.9),Inches(2.4),BLUE_BG)
        txb(sl,badge,bx+Inches(.12),by+Inches(.1),Inches(3.65),Inches(.35),size=13,bold=True,color=BLUE_L)
        title,desc=body.split('\n',1)
        txb(sl,title,bx+Inches(.12),by+Inches(.52),Inches(3.65),Inches(.35),size=13,bold=True,color=CHARCOAL)
        txb(sl,desc, bx+Inches(.12),by+Inches(.92),Inches(3.65),Inches(1.3),size=11,color=GRAY)

    # ── 05 六大指標總覽 ────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 指標架構',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'六大評估指標總覽',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    indicators=[
        ('01','📡','資訊通信','ICT Infrastructure',      C_ICT, '#eff6ff'),
        ('02','🛡️','安全防災','Safety & Security',        C_SEC, '#fff1f2'),
        ('03','🌡️','健康舒適','Health & Comfort',         C_HLTH,'#f0fdf4'),
        ('04','⚡','節能管理','Energy Management',         C_NRG, '#fffbeb'),
        ('05','🔗','設備整合','Systems Integration',       C_INT, '#f5f3ff'),
        ('06','📱','貼心便利','Smart Convenience',         C_CVN, '#ecfeff'),
    ]
    for i,(num,icon,name,en,fc,bg) in enumerate(indicators):
        col=i%3; row=i//2 if i<3 else (i-3)//3+1
        # arrange 2 rows of 3
        row2=i//3
        bx=Inches(1+col*4.1); by=Inches(2.0+row2*2.5)
        rect(sl,bx,by,Inches(3.9),Inches(2.3),bg)
        txb(sl,num, bx+Inches(.12),by+Inches(.12),Inches(.5),Inches(.38),size=14,bold=True,color=fc)
        txb(sl,icon,bx+Inches(.12),by+Inches(.55),Inches(3.65),Inches(.5),size=24)
        txb(sl,name,bx+Inches(.12),by+Inches(1.1),Inches(3.65),Inches(.4),size=16,bold=True,color=fc)
        txb(sl,en,  bx+Inches(.12),by+Inches(1.52),Inches(3.65),Inches(.3),size=11,color=GRAY,italic=True)

    # ── 06 配分比重 ────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 配分比重',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'六大指標配分比重',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    weights=[('資訊通信',20,C_ICT),('安全防災',20,C_SEC),('健康舒適',15,C_HLTH),
             ('節能管理',20,C_NRG),('設備整合',15,C_INT),('貼心便利',10,C_CVN)]
    for i,(name,pct,fc) in enumerate(weights):
        by=Inches(2.0+i*.76)
        txb(sl,name,Inches(1),by,Inches(2.2),Inches(.6),size=13,color=fc)
        bar_w=Inches(6.5*pct/100)
        rect(sl,Inches(3.3),by+Inches(.12),bar_w,Inches(.38),fc)
        txb(sl,f'{pct}%',Inches(3.3)+bar_w+Inches(.1),by+Inches(.1),Inches(.8),Inches(.42),size=13,bold=True,color=fc)
    rect(sl,Inches(7.8),Inches(2.0),Inches(5),Inches(2.2),BLUE_BG)
    txb(sl,'評分架構說明',Inches(7.95),Inches(2.1),Inches(4.7),Inches(.35),size=13,bold=True,color=BLUE_L)
    txb_lines(sl,['各指標分為必要項目與選擇（加分）項目',
                   '必要項目未達標，該指標不計分',
                   '加分項目每項可累積得分，達上限後不再計算',
                   '2024 版各指標均設有「AI 智慧化」加分子項'],
        Inches(7.95),Inches(2.5),Inches(4.7),Inches(1.5),size=12,color=CHARCOAL,bullet=True)
    rect(sl,Inches(7.8),Inches(4.35),Inches(5),Inches(1.9),BLUE_BG)
    txb(sl,'適用建築類型',Inches(7.95),Inches(4.45),Inches(4.7),Inches(.35),size=13,bold=True,color=BLUE_L)
    txb_lines(sl,['辦公大樓','商業建築','住宅大廈','學校教育','醫療院所','廠辦建築'],
        Inches(7.95),Inches(4.85),Inches(4.7),Inches(1.1),size=12,color=CHARCOAL,bullet=True)

    # ── 07-12 Six indicators ───────────────────
    ind_data=[
        ('01','資訊通信','ICT Infrastructure',C_ICT,'#eff6ff',
         ['結構化佈線：Cat.6A 以上線纜或光纖主幹',
          '無線網路覆蓋：全棟 Wi-Fi 6／6E 無縫覆蓋',
          '機房設施：UPS 不斷電、溫控、獨立門禁',
          '資訊安全基礎：防火牆、網路分區隔離'],
         ['5G 小基站室內部署 [2024新增]',
          '零信任（Zero Trust）資安架構 [2024強化]',
          'IDS/IPS 入侵偵測防禦系統',
          'N+1 冗餘備援機制',
          'AI 網路流量異常偵測 [2024 AI加分]'],
         '2024 重點：5G × Zero Trust',
         '5G 小基站讓建築內部具備毫秒級低延遲連線。\n零信任架構「永不信任、持續驗證」，大幅提升建築資訊系統的資安等級。'),
        ('02','安全防災','Safety & Security',C_SEC,'#fff1f2',
         ['智慧門禁：IC 卡、藍牙、生物辨識（至少一種）',
          'CCTV 監控：主要公共區域全覆蓋',
          '火警系統：智慧型煙霧偵測、自動灑水連動',
          '緊急廣播：數位多區域廣播系統'],
         ['AI 影像行為分析（異常偵測）[2024新增]',
          '人臉辨識門禁 × 訪客預約整合',
          '熱感應周界防護',
          '地震感測器 × 自動斷電保護',
          'AI 火焰早期偵測 [2024 AI加分]'],
         '2024 重點：AI 影像分析',
         'AI 影像行為分析能自動辨識人員聚集、遺留物、翻越圍欄等異常行為，\n即時發出警報，從被動監控進化為主動預防的智慧安防系統。'),
        ('03','健康舒適','Health & Comfort',C_HLTH,'#f0fdf4',
         ['室內空氣品質（IAQ）監測：CO₂ 感測器裝設',
          '空調自動控制：依設定溫度自動調節',
          '照明自動控制：人體感應或時序控制'],
         ['PM2.5、TVOC 即時監測 + 顯示面板',
          'AI 預測性空調控制（減少啟停震盪）[2024 AI加分]',
          '智慧照明：晝光感測 × 色溫自動調節',
          '給水品質智慧監測（餘氯、濁度）',
          '電梯目的地控制（DCS）智慧調度'],
         '2024 重點：AI 預測性空調',
         'AI 預測性控制透過分析室外氣溫、建築熱慣性、使用人數等資料，\n提前調整空調輸出，同步達到舒適度提升與能耗降低的雙重效益。'),
        ('04','節能管理','Energy Management',C_NRG,'#fffbeb',
         ['能源管理系統（EMS）：全棟用電即時監控儀表',
          '分項計量：空調、照明、動力各自獨立計量',
          '智慧電表（AMI）：30 分鐘間隔自動讀表'],
         ['太陽光電 + ESS 儲能整合 [2024強化]',
          'AI 電力需量預測與自動控制 [2024 AI加分]',
          '建築碳足跡即時監控儀表板 [2024新增]',
          '電動車充電樁智慧管理',
          '需量反應（DR）自動化機制'],
         '2024 重點：碳足跡監控',
         '2024 版新增建築碳足跡即時監控加分項：\n系統依能源耗用量自動換算 CO₂ 排放，提供月度碳排告報告，協助建築物進行碳盤查。'),
        ('05','設備整合','Systems Integration',C_INT,'#f5f3ff',
         ['建築自動化系統（BAS / BMS）：跨系統整合平台',
          '中央監控室：統一操作介面（至少三系統整合）',
          '開放通訊協議：BACnet、Modbus 或 MQTT'],
         ['BIM 數位孿生整合（3D 即時管理）[2024新增]',
          'AI 設備壽命預測性維護 [2024 AI加分]',
          '雲端遠端監控平台',
          'IoT 閘道器 × 感測器網路擴充',
          '數位維運工單系統（CMMS）'],
         '2024 重點：BIM 數位孿生',
         '數位孿生將實體建築 1:1 映射為 3D 數位模型，\n每個設備的即時狀態直接顯示在模型上，大幅提升維護效率並減少停機損失。'),
        ('06','貼心便利','Smart Convenience',C_CVN,'#ecfeff',
         ['智慧停車：車牌辨識 + 剩餘車位即時顯示',
          '訪客管理：電子登記 + 訪客證件核驗',
          '資訊服務：大廳資訊看板 / 導覽系統'],
         ['物業管理 App（報修、繳費、公告）[2024強化]',
          '智慧信箱：包裹自動通知 + 自取',
          '共用設施線上預約系統',
          '智慧停車自動繳費（ETC / 行動支付）',
          'IoT 智慧家電整合控制 [2024新增]'],
         '2024 重點：物業 App 整合',
         '2024 版大幅提升行動 App 整合評分比重，\n要求物業 App 涵蓋：報修申請、社區公告、訪客 QR Code 邀請、共用設施預約、繳費等。'),
    ]
    for num,name,en,fc,bg,must,opt,hd,hd_desc in ind_data:
        sl = blank_slide(prs)
        fill_bg(sl,'#fafaf9')
        txb(sl,f'{num} {name} ／ {en}',Inches(1),Inches(.5),Inches(8),Inches(.4),size=12,color=fc)
        accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),fc)
        txb(sl,f'{name} — ',Inches(1),Inches(1.07),Inches(7),Inches(.6),size=26,bold=True,color=CHARCOAL)
        txb(sl,en,Inches(1),Inches(1.67),Inches(7),Inches(.4),size=14,color=GRAY,italic=True)
        txb(sl,'必要項目',Inches(1),Inches(2.2),Inches(6),Inches(.35),size=13,bold=True,color=fc)
        txb_lines(sl,must,Inches(1),Inches(2.6),Inches(6),Inches(1.6),size=12,color=CHARCOAL,bullet=True)
        txb(sl,'選擇（加分）項目',Inches(1),Inches(4.3),Inches(6),Inches(.35),size=13,bold=True,color=fc)
        txb_lines(sl,opt,Inches(1),Inches(4.7),Inches(6),Inches(1.8),size=11,color=CHARCOAL,bullet=True)
        rect(sl,Inches(7.5),Inches(2.0),Inches(5.5),Inches(2.2),bg)
        txb(sl,hd,Inches(7.65),Inches(2.1),Inches(5.2),Inches(.35),size=13,bold=True,color=fc)
        txb(sl,hd_desc,Inches(7.65),Inches(2.5),Inches(5.2),Inches(1.5),size=12,color=CHARCOAL)

    # ── 13 四個認證等級 ────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 認證等級',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'四個認證等級',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    levels=[
        ('🏅','合格級','Qualified','基礎門檻','#6b7280',
         '各指標必要項目全數通過，選擇項目達最低得分門檻，具備完整智慧建築基礎功能。'),
        ('🥈','銀　級','Silver',  '優良表現','#64748b',
         '在合格基礎上，多項加分指標有良好表現，系統整合度高，具備智慧化升級潛力。'),
        ('🥇','黃金級','Gold',    '卓越表現','#b45309',
         '多數指標達高分，具備 AI 智慧化、跨系統整合及主動式管理能力，為市場標竿案例。'),
        ('💎','鑽石級','Diamond', '頂尖典範','#1e40af',
         '六大指標均達頂尖水準，涵蓋 AI 分析、數位孿生、5G 及碳足跡監控，代表台灣最高智慧建築標準。'),
    ]
    for i,(icon,name,en,score,fc,desc) in enumerate(levels):
        bx=Inches(1+i*3.1)
        rect(sl,bx,Inches(1.85),Inches(2.9),Inches(3.5),'#f9fafb')
        txb(sl,icon,bx,Inches(1.95),Inches(2.9),Inches(.6),size=24,align=PP_ALIGN.CENTER)
        txb(sl,name,bx,Inches(2.65),Inches(2.9),Inches(.4),size=15,bold=True,color=fc,align=PP_ALIGN.CENTER)
        txb(sl,en,  bx,Inches(3.05),Inches(2.9),Inches(.3),size=11,color=GRAY,align=PP_ALIGN.CENTER)
        txb(sl,score,bx,Inches(3.35),Inches(2.9),Inches(.3),size=12,bold=True,color=fc,align=PP_ALIGN.CENTER)
        txb(sl,desc,bx+Inches(.1),Inches(3.72),Inches(2.7),Inches(1.4),size=10,color=GRAY)
    rect(sl,Inches(1),Inches(5.5),Inches(5.8),Inches(1.65),BLUE_BG)
    txb(sl,'2024 版鑽石級新要求',Inches(1.15),Inches(5.6),Inches(5.5),Inches(.35),size=13,bold=True,color=BLUE_L)
    txb_lines(sl,['至少導入 2 項 AI 智慧化加分子項','完成建築碳足跡監控系統建置','具備 BIM 數位孿生或 5G 基礎設施'],
        Inches(1.15),Inches(6.0),Inches(5.5),Inches(.9),size=12,color=CHARCOAL,bullet=True)
    rect(sl,Inches(7.1),Inches(5.5),Inches(5.8),Inches(1.65),BLUE_BG)
    txb(sl,'評定有效期與展延',Inches(7.25),Inches(5.6),Inches(5.5),Inches(.35),size=13,bold=True,color=BLUE_L)
    txb_lines(sl,['標章有效期：5 年','到期前 6 個月可申請重新評定或展延','系統升級後可申請等級提升評定'],
        Inches(7.25),Inches(6.0),Inches(5.5),Inches(.9),size=12,color=CHARCOAL,bullet=True)

    # ── 14 申請流程 ────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 申請流程',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'申請流程 — 四大階段',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    steps=[
        ('STEP 01','📋','設計階段規劃申請','確認智慧化系統設計規格，向建研所提出「候選智慧建築證書」申請，確立評估方向'),
        ('STEP 02','🏗️','施工階段文件蒐集','設備採購規格書、施工紀錄、系統整合測試報告、資安掃描報告等文件備齊'),
        ('STEP 03','🔍','竣工後現場勘查','委託評定機構進行書面審查＋現場設備功能測試，驗證各指標實際運作狀況'),
        ('STEP 04','🏆','核發標章公開揭示','通過評定後核發正式標章，建築物公共區域需揭示標章，並登錄至建研所資料庫'),
    ]
    for i,(num,icon,title,desc) in enumerate(steps):
        bx=Inches(1+i*3.1)
        rect(sl,bx,Inches(2.0),Inches(2.9),Inches(2.8),BLUE_BG)
        txb(sl,num, bx+Inches(.12),Inches(2.1),Inches(2.65),Inches(.3),size=11,color=BLUE_L)
        txb(sl,icon,bx+Inches(.12),Inches(2.45),Inches(2.65),Inches(.45),size=22)
        txb(sl,title,bx+Inches(.12),Inches(2.97),Inches(2.65),Inches(.5),size=13,bold=True,color=CHARCOAL)
        txb(sl,desc, bx+Inches(.12),Inches(3.52),Inches(2.65),Inches(1.1),size=10,color=GRAY)
        if i<3:
            txb(sl,'→',Inches(3.9+i*3.1),Inches(3.2),Inches(.3),Inches(.4),size=16,color=GRAY)
    rect(sl,Inches(1),Inches(5.0),Inches(5.8),Inches(2.1),'#f9fafb')
    txb(sl,'申請應備文件',Inches(1.15),Inches(5.1),Inches(5.5),Inches(.35),size=13,bold=True,color=CHARCOAL)
    txb_lines(sl,['智慧化系統設計說明書','各子系統設備規格書及型錄',
                   '系統整合測試報告（FAT/SAT）','資訊安全評估報告（2024版強制）','竣工平面圖 × 系統架構圖'],
        Inches(1.15),Inches(5.5),Inches(5.5),Inches(1.4),size=12,color=CHARCOAL,bullet=True)
    rect(sl,Inches(7.1),Inches(5.0),Inches(5.8),Inches(2.1),BLUE_BG)
    txb(sl,'評定機構',Inches(7.25),Inches(5.1),Inches(5.5),Inches(.35),size=13,bold=True,color=BLUE_L)
    txb(sl,'申請人可自行向內政部建築研究所申請，或委託建研所認可之評定機構代為辦理。\n評定費用依建築規模與評估指標數量而定。',
        Inches(7.25),Inches(5.5),Inches(5.5),Inches(1.4),size=12,color=CHARCOAL)

    # ── 15 快速回顧表 ──────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl,'#fafaf9')
    txb(sl,'智慧建築 ／ 重點回顧',Inches(1),Inches(.5),Inches(5),Inches(.4),size=12,color=BLUE_L)
    accent_bar(sl,Inches(1),Inches(.93),Inches(1.2),BLUE_L)
    txb(sl,'六大指標 × 2024 更新一覽',Inches(1),Inches(1.07),Inches(11),Inches(.6),size=28,bold=True,color=CHARCOAL)
    # header
    rect(sl,Inches(1),Inches(1.85),Inches(2.5),Inches(.42),BLUE)
    txb(sl,'指標',Inches(1.05),Inches(1.9),Inches(2.4),Inches(.35),size=12,bold=True,color=WHITE)
    rect(sl,Inches(3.55),Inches(1.85),Inches(3.0),Inches(.42),BLUE)
    txb(sl,'核心功能',Inches(3.6),Inches(1.9),Inches(2.9),Inches(.35),size=12,bold=True,color=WHITE)
    rect(sl,Inches(6.6),Inches(1.85),Inches(3.0),Inches(.42),BLUE)
    txb(sl,'必要項目',Inches(6.65),Inches(1.9),Inches(2.9),Inches(.35),size=12,bold=True,color=WHITE)
    rect(sl,Inches(9.65),Inches(1.85),Inches(3.0),Inches(.42),BLUE)
    txb(sl,'2024 新重點',Inches(9.7),Inches(1.9),Inches(2.9),Inches(.35),size=12,bold=True,color=WHITE)
    table_data=[
        (C_ICT,'① 資訊通信','ICT 基礎 × 網路覆蓋','Cat.6A 佈線、Wi-Fi 全棟','5G小基站、Zero Trust'),
        (C_SEC,'② 安全防災','門禁 × 監控 × 火警','門禁、CCTV、火警、廣播','AI 影像行為分析'),
        (C_HLTH,'③ 健康舒適','IAQ × 空調 × 照明','CO₂ 監測、空調自控','AI 預測性空調控制'),
        (C_NRG,'④ 節能管理','EMS × 計量 × 再生能源','EMS、分項計量、智慧電表','碳足跡監控儀表板'),
        (C_INT,'⑤ 設備整合','BAS/BMS × 跨系統整合','BMS、開放協議、中控室','BIM 數位孿生'),
        (C_CVN,'⑥ 貼心便利','停車 × 訪客 × 行動服務','智慧停車、訪客管理','物業 App × IoT 家電整合'),
    ]
    for i,(fc,ind,func,must,new) in enumerate(table_data):
        bg='#f9fafb' if i%2==0 else WHITE
        y=Inches(2.35+i*.76)
        rect(sl,Inches(1),y,Inches(2.5),Inches(.7),bg)
        txb(sl,ind,Inches(1.05),y+Pt(4),Inches(2.4),Inches(.62),size=12,bold=True,color=fc)
        for j,(col_x,col_w,txt) in enumerate([(3.55,3.0,func),(6.6,3.0,must),(9.65,3.0,new)]):
            rect(sl,Inches(col_x),y,Inches(col_w),Inches(.7),bg)
            txb(sl,txt,Inches(col_x+.08),y+Pt(4),Inches(col_w-.16),Inches(.62),size=11,color=CHARCOAL)

    # ── 16 結語 ────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, BLUE_DARK)
    txb(sl,'智慧建築標章 2024 版',Inches(1),Inches(.8),Inches(10),Inches(.4),size=13,color=BLUE_LG)
    txb(sl,'從 ICT 基礎到 AI 數位孿生，',Inches(1),Inches(1.4),Inches(11),Inches(.7),size=34,bold=True,color=WHITE)
    txb(sl,'智慧建築不只是科技',Inches(1),Inches(2.1),Inches(11),Inches(.7),size=34,bold=True,color=BLUE_LG)
    txb(sl,'而是對使用者生活品質的完整承諾',Inches(1),Inches(2.8),Inches(11),Inches(.6),size=22,italic=True,color=WHITE)
    rect(sl,Inches(1),Inches(3.55),Inches(.7),Pt(4),BLUE_LG)
    # indicator dots
    for i,(name,fc) in enumerate([('資訊通信',C_ICT),('安全防災',C_SEC),('健康舒適',C_HLTH),
                                    ('節能管理',C_NRG),('設備整合',C_INT),('貼心便利',C_CVN)]):
        bx=Inches(1+i*2.0)
        rect(sl,bx,Inches(4.0),Inches(.18),Inches(.18),fc)
        txb(sl,name,bx+Inches(.25),Inches(3.96),Inches(1.7),Inches(.28),size=12,color='#94a3b8')
    txb(sl,'資料來源：內政部建築研究所 智慧建築標章評估手冊 2024 版',
        Inches(1),Inches(6.8),Inches(11),Inches(.35),size=11,color='#334155',italic=True)

    prs.save(path)
    print(f'Saved {path}')

build_ib('/home/user/jieceng-web/presentation-IB.pptx')
print('IB done')
